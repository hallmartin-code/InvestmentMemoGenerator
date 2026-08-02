"""Generate synthetic pitch decks (PDF + PPTX) used to exercise the parser.

Run: python tests/make_sample_deck.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import landscape, letter
from reportlab.pdfgen import canvas

HERE = Path(__file__).resolve().parent

SLIDES = [
    ("Northwind Grid", ["Battery-free load shifting for commercial buildings",
                        "Seed round — 2026"]),
    ("The Problem", [
        "Commercial buildings pay demand charges on their single highest 15-minute peak.",
        "Those charges are 30-70% of a typical monthly bill.",
        "Batteries fix it but cost $400k+ per site and take 14 months to install.",
    ]),
    ("Our Solution", [
        "Northwind Grid is a control layer that pre-cools and staggers HVAC and",
        "refrigeration loads to shave the peak without any hardware on site.",
        "Installs in 3 days over existing BMS protocols (BACnet, Modbus).",
    ]),
    ("How It Works", [
        "Forecast the building's load curve 24h ahead from meter and weather data.",
        "Solve for a dispatch schedule that keeps every zone inside its comfort band.",
        "Push setpoints to the BMS. Fall back to native control on any fault.",
    ]),
    ("Traction", [
        "41 buildings live across 6 customers.",
        "$780k ARR, growing 14% month over month for the last 7 months.",
        "Average verified demand-charge reduction: 22%.",
        "Net revenue retention 131%. Two customers expanded portfolio-wide.",
    ]),
    ("Business Model", [
        "Shared savings: we take 35% of verified demand-charge reduction,",
        "measured against an IPMVP Option C baseline.",
        "No hardware, no capex, no upfront fee. Contracts are 3 years.",
    ]),
    ("Market", [
        "5.9M commercial buildings in the US.",
        "1.1M are over 25,000 sq ft with a BMS we can reach today (SAM).",
        "At $9k average annual revenue per building that is a $9.9B SAM.",
        "Source: EIA CBECS 2018, internal pricing model.",
    ]),
    ("Team", [
        "Priya Raman, CEO. 9 years at Enel X leading C&I demand response.",
        "Tom Achebe, CTO. Built the dispatch optimizer at Voltus.",
        "Four engineers. No commercial hires yet.",
    ]),
    ("Competition", [
        "Batteries (Stem, Sunrun): higher savings, 100x the cost and install time.",
        "Legacy demand response (Enel X, CPower): event-driven, utility-triggered only.",
        "BMS vendors (JCI, Siemens): bundled, slow to ship, weak forecasting.",
    ]),
    ("The Ask", [
        "Raising $6M Seed.",
        "60% engineering (protocol coverage, M&V automation).",
        "25% go-to-market (first 3 AEs, 1 solutions engineer).",
        "15% working capital.",
        "24 months of runway to $6M ARR.",
    ]),
]


def build_pdf(path):
    c = canvas.Canvas(str(path), pagesize=landscape(letter))
    width, height = landscape(letter)
    for title, lines in SLIDES:
        c.setFont("Helvetica-Bold", 28)
        c.drawString(72, height - 110, title)
        c.setFont("Helvetica", 16)
        y = height - 170
        for line in lines:
            c.drawString(72, y, line)
            y -= 28
        c.showPage()
    c.save()


def build_pptx(path):
    presentation = Presentation()
    blank = presentation.slide_layouts[6]
    for title, lines in SLIDES:
        slide = presentation.slides.add_slide(blank)
        title_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(0.5), Inches(9), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].runs[0].font.size = Pt(32)
        title_frame.paragraphs[0].runs[0].font.bold = True

        body_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.7), Inches(9), Inches(4)
        )
        body_frame = body_box.text_frame
        body_frame.word_wrap = True
        for index, line in enumerate(lines):
            paragraph = (
                body_frame.paragraphs[0] if index == 0 else body_frame.add_paragraph()
            )
            run = paragraph.add_run()
            run.text = line
            run.font.size = Pt(18)
    presentation.save(str(path))


if __name__ == "__main__":
    build_pdf(HERE / "sample_deck.pdf")
    build_pptx(HERE / "sample_deck.pptx")
    print(f"Wrote {HERE / 'sample_deck.pdf'}")
    print(f"Wrote {HERE / 'sample_deck.pptx'}")
