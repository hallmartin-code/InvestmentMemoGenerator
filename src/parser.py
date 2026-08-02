"""Deck ingestion: turn a PDF or PPTX pitch deck into a single raw text string."""

from pathlib import Path

PAGE_BREAK = "\n--- PAGE BREAK ---\n"

SUPPORTED_EXTENSIONS = (".pdf", ".pptx", ".docx")


class UnsupportedFormatError(Exception):
    """Raised when the input file is not a .pdf or .pptx."""


class DeckParseError(Exception):
    """Raised when a deck cannot be read or contains no extractable text."""


def parse_deck(path):
    """Parse a pitch deck into raw text.

    Slide/page boundaries are marked with PAGE_BREAK so the extractor can
    reason about them.
    """
    path = Path(path)
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        raw_text = _parse_pdf(path)
    elif suffix == ".pptx":
        raw_text = _parse_pptx(path)
    elif suffix == ".docx":
        raw_text = _parse_docx(path)
    else:
        raise UnsupportedFormatError(suffix)

    if not raw_text.strip():
        raise DeckParseError(
            f"No extractable text found in {path.name}. "
            "The deck may be image-only (scanned or exported as flat images)."
        )
    return raw_text


def _parse_pdf(path):
    import pdfplumber

    pages = []
    try:
        with pdfplumber.open(str(path)) as pdf:
            for page in pdf.pages:
                pages.append(page.extract_text() or "")
    except Exception as exc:
        raise DeckParseError(f"Could not read PDF {path.name}: {exc}") from exc

    return PAGE_BREAK.join(pages)


def _parse_pptx(path):
    from pptx import Presentation

    try:
        presentation = Presentation(str(path))
    except Exception as exc:
        raise DeckParseError(f"Could not read PPTX {path.name}: {exc}") from exc

    slides = []
    for index, slide in enumerate(presentation.slides, start=1):
        lines = [f"[Slide {index}]"]

        title_shape = slide.shapes.title
        title_text = ""
        if title_shape is not None and title_shape.has_text_frame:
            title_text = _shape_text(title_shape)
            if title_text:
                lines.append(f"TITLE: {title_text}")

        for shape in slide.shapes:
            if shape is title_shape:
                continue
            text = _shape_text(shape)
            if text:
                lines.append(text)

        slides.append("\n".join(lines))

    return PAGE_BREAK.join(slides)


def _parse_docx(path):
    """Extract text from a Word deck or memo.

    Explicit page breaks become PAGE_BREAK markers; headings are prefixed so
    the extractor can see document structure.
    """
    import docx
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    try:
        document = docx.Document(str(path))
    except Exception as exc:
        raise DeckParseError(f"Could not read DOCX {path.name}: {exc}") from exc

    parts = []
    for block in _iter_docx_blocks(document, Paragraph, Table):
        if isinstance(block, Paragraph):
            text = block.text.strip()
            if not text:
                continue
            style = (block.style.name or "") if block.style is not None else ""
            parts.append(f"HEADING: {text}" if style.startswith("Heading") else text)
            if _has_page_break(block):
                parts.append(PAGE_BREAK.strip())
        else:
            for row in block.rows:
                cells = [cell.text.strip() for cell in row.cells]
                row_text = " | ".join(c for c in cells if c)
                if row_text:
                    parts.append(row_text)

    return "\n".join(parts)


def _iter_docx_blocks(document, Paragraph, Table):
    """Yield paragraphs and tables in document order."""
    from docx.oxml.ns import qn

    body = document.element.body
    for child in body.iterchildren():
        if child.tag == qn("w:p"):
            yield Paragraph(child, document)
        elif child.tag == qn("w:tbl"):
            yield Table(child, document)


def _has_page_break(paragraph):
    from docx.oxml.ns import qn

    for run in paragraph.runs:
        for br in run._element.findall(qn("w:br")):
            if br.get(qn("w:type")) == "page":
                return True
    return False


def _shape_text(shape):
    """Concatenate text from all runs in a shape, preserving paragraph breaks.

    Recurses into grouped shapes and reads table cells, which is where a lot
    of pitch-deck copy lives.
    """
    parts = []

    if getattr(shape, "has_text_frame", False):
        for paragraph in shape.text_frame.paragraphs:
            line = "".join(run.text for run in paragraph.runs).strip()
            if line:
                parts.append(line)

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            row_text = " | ".join(c for c in cells if c)
            if row_text:
                parts.append(row_text)

    # shape_type 6 == GROUP; guard with getattr so non-group shapes are skipped.
    if getattr(shape, "shapes", None) is not None:
        for child in shape.shapes:
            child_text = _shape_text(child)
            if child_text:
                parts.append(child_text)

    return "\n".join(parts)
